from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from langchain_pptx2markdown import PPTX2MarkdownLoader


class EndToEndTests(unittest.TestCase):
    def test_generated_pptx_loads_without_persistent_work_directories(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "한글 deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            text_box.text = "Hello from PowerPoint"
            presentation.save(source)

            documents = PPTX2MarkdownLoader(source).load()

            self.assertEqual(len(documents), 1)
            self.assertIn("Hello from PowerPoint", documents[0].page_content)
            self.assertEqual(documents[0].metadata["slide_number"], 1)
            self.assertEqual(documents[0].metadata["file_name"], source.name)
            self.assertFalse((root / "output").exists())
            self.assertFalse((root / ".pptx2markdown").exists())


if __name__ == "__main__":
    unittest.main()
